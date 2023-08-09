
from flask import Flask, request
# from document_generator import generate_docx, generate_pptx
import json
from docx import Document
import base64
from pptx import Presentation
from pptx.util import Inches, Pt

app = Flask(__name__)


@app.route('/', methods=['GET'])
def index():
    return "Flask API works successfully"

@app.route('/generate_docx', methods=['POST'])
def generate_docx():
    doc = Document()
    text = request.json['text']
    doc.add_paragraph(text)
    bytesData = doc.save('output.docx')
    with open('output.docx', 'rb') as file:
        result = base64.b64encode(file.read())
    return result.decode('utf-8')


@app.route('/generate_ppt', methods=['POST'])
def generate_ppt():
    print('reacheddd')
    text = request.json['text']
    # Creating presentation object
    ppt = Presentation()
    # To create blank slide layout
    # We have to use 6 as an argument
    # of slide_layouts  
    blank_slide_layout = ppt.slide_layouts[6] 
    
    # Attaching slide obj to slide
    slide = ppt.slides.add_slide(blank_slide_layout)
    
    # For adjusting the  Margins in inches 
    left = top = width = height = Inches(1) 
    
    # creating textBox
    txBox = slide.shapes.add_textbox(left, top,
                                    width, height)
    # creating textFrames
    tf = txBox.text_frame
    # adding Paragraphs
    p = tf.add_paragraph() 
    # adding text
    p.text = text
    
    # font 
    p.font.bold = True
    p.font.italic = True

    # save file
    ppt.save('test_1.pptx')
    with open('test_1.pptx', 'rb') as file:
        result = base64.b64encode(file.read())
    return result.decode('utf-8')


# main driver function
if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000, debug=True)
