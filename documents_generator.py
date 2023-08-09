# from docx import Document
# from pptx import Presentation
# from io import BytesIO
# import sys
# import base64


# def generate_docx(text):
#     doc = Document()
#     doc.add_paragraph(text)
#     doc_bytes = BytesIO()
#     doc.save(doc_bytes)
#     doc_bytes.seek(0)
#     return base64.b64encode(doc_bytes.read()).decode("utf-8")


# def generate_pptx(text):
#     prs = Presentation()
#     slide = prs.slides.add_slide(prs.slide_layouts[5])
#     slide.shapes.title.text = text
#     ppt_bytes = BytesIO()
#     prs.save(ppt_bytes)
#     ppt_bytes.seek(0)
#     return base64.b64encode(ppt_bytes.read()).decode("utf-8")


# if __name__ == "__main__":
#     if len(sys.argv) == 2:
#         text = sys.argv[1]
#         format_type = sys.argv[2]
#         if format_type == "docx":
#             result = generate_docx(text)
#             print(result)
#         elif format_type == "pptx":
#             result = generate_pptx(text)
#             print(result)
#         else:
#             print("Invalid format type. Supported formats: docx, pptx")
#     else:
#         print("Usage: python document_generator.py <text> <format_type>")


# document_generator.py

import sys
from docx import Document
from pptx import Presentation
import base64

def generate_docx(text):
    doc = Document()
    doc.add_paragraph(text)
    result = base64.b64encode(doc.save_to_bytes())
    return result.decode('utf-8')

def generate_pptx(text):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    text_box = slide.shapes.add_textbox(left=0, top=0, width=prs.slide_width, height=prs.slide_height)
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    result = base64.b64encode(prs.save_to_bytes())
    return result.decode('utf-8')

if __name__ == '__main__':
    command = sys.argv[1]
    text = sys.argv[2]
    if command == 'docx':
        result = generate_docx(text)
        print(result)
    elif command == 'pptx':
        result = generate_pptx(text)
        print(result)

